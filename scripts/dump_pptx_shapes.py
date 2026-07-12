#!/usr/bin/env python3
"""Diagnostic dump of every shape's text + position in a PPTX file.

Not part of the app -- a throwaway helper to inspect the geometry of the
old manual-format PPTX before writing the real importer, since shapes have
no meaningful names (all "TextBox N"/"Grupo N") and column assignment has
to be reconstructed from (x, y) position instead.

Usage: python dump_pptx_shapes.py archivo.pptx
"""
import sys
from pptx import Presentation
from pptx.util import Emu


def emu_to_in(v):
    return round(Emu(v).inches, 2) if v is not None else None


def walk(shapes, depth=0):
    for shape in shapes:
        left, top = emu_to_in(shape.left), emu_to_in(shape.top)
        width, height = emu_to_in(shape.width), emu_to_in(shape.height)
        text = ""
        if shape.has_text_frame:
            text = shape.text_frame.text.replace("\n", " \\n ").strip()
        prefix = "  " * depth
        print(f"{prefix}[{shape.shape_type}] pos=({left},{top}) size=({width},{height}) text={text!r}")
        if shape.shape_type == 6:  # GROUP
            walk(shape.shapes, depth + 1)


def main():
    path = sys.argv[1]
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, 1):
        print(f"\n=================== SLIDE {i} ===================")
        walk(slide.shapes)


if __name__ == "__main__":
    main()
