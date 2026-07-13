#!/usr/bin/env python3
"""Import incidents from the old manual-format PPTX (pre-dating this app)
into a JSON payload matching POST /api/reports.

The old deck has no per-incident type information this app tracks
(severity, brands, ministry/platform/external-origin flags, action points) -- those are
always defaulted and MUST be reviewed by hand after import. Everything
else (title, date, duration, ticket(s), the Impacto/Causa/Solución text,
and the metrics lines) is reconstructed from each shape's position on the
slide, since the shapes themselves have no meaningful names (see
dump_pptx_shapes.py, used to reverse-engineer this layout).

Usage:
  python import_pptx.py archivo.pptx -o salida.json
  python import_pptx.py archivo.pptx -o salida.json --year 2026 --week 28 --range "06-12 julio 2026" --dept "Customer & Service Operations"
  python import_pptx.py archivo.pptx --post-url http://localhost:8000   # crea el informe directamente via la API

If --year/--week aren't given, they're parsed from the filename
(pattern "2026W28"). --range defaults to the min-max dates found in the
incidents themselves if omitted.
"""
import argparse
import json
import re
import sys
import urllib.request

from pptx import Presentation

# Windows consoles default to cp1252, which can't encode the emoji/accented
# characters this script prints -- force UTF-8 regardless of the terminal.
sys.stdout.reconfigure(encoding="utf-8")
sys.stderr.reconfigure(encoding="utf-8")

# Real ticket codes (e.g. "2607C80395") always contain a digit; the column
# label words below are pure alphabetic and would otherwise false-match a
# plain [A-Za-z0-9]{6,}.
TICKET_RE = re.compile(r"(?=.*\d)[A-Za-z0-9]{6,}")

# Column/field label textboxes -- excluded wherever we're looking for a
# *value* near a label's coordinates, since on at least one slide the label
# and its value share the same local x (see nearest()'s tie-breaking note).
LABEL_WORDS = ("Impacto", "Causa", "Solución", "ID", "Fecha", "Duración")


def slide_group(title_text):
    """Map the slide's header text to this app's `group` field.
    Falls back to a literal "RED <lo que ponga>" so nothing gets silently
    dropped, but flags it since it's not one of the known conventions."""
    t = title_text.replace("Incidencias RED", "").strip(" ()")
    tl = t.lower()
    if "5000" in tl or "5.000" in tl:
        return "RED >5.000 clientes", True
    if "relevantes" in tl or "climatolog" in tl or "escalados" in tl:
        return "Otras RED", True
    if "b2b" in tl:
        return "RED B2B", False  # not an existing convention -- flagged
    return f"RED ({t})" if t else "Otras RED", False


def emu_in(v):
    from pptx.util import Emu
    return round(Emu(v).inches, 2) if v is not None else None


def collect_textboxes(shapes, out):
    """Recursively collect (x, y, text) for every text-bearing shape.

    Deliberately does NOT accumulate parent-group offsets: a PPTX group's
    children are positioned in the group's own child coordinate space
    (chOff/chExt), which maps to the group's outer bounding box via an
    affine transform (offset *and* scale) -- naively adding the parent's
    left/top at each nesting level double-counts that and produces
    nonsense coordinates. What we actually need is simpler: every
    incident group on this deck repeats the *same* internal template, so
    each field's raw (non-accumulated) shape.left/top is consistent
    incident to incident (e.g. the title textbox is always at local
    x=0.53, y=0.84) -- confirmed against dump_pptx_shapes.py's output.
    That local consistency is all the classification below relies on."""
    for shape in shapes:
        x = emu_in(shape.left) or 0
        y = emu_in(shape.top) or 0
        if shape.shape_type == 6:  # GROUP
            collect_textboxes(shape.shapes, out)
        elif shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                out.append((x, y, text))


def nearest(items, target_x, tolerance=0.6):
    best = None
    best_dist = tolerance
    for x, y, text in items:
        d = abs(x - target_x)
        if d < best_dist:
            best_dist = d
            best = text
    return best


def metrics_to_pipe_format(text):
    """'Clientes de FTTH: 94K' -> 'Clientes de FTTH | 94K', one per line
    (matches this app's metricsArr() "label | value" convention)."""
    lines = []
    for line in text.split("\n"):
        line = line.strip()
        if not line:
            continue
        if ":" in line:
            label, _, value = line.partition(":")
            lines.append(f"{label.strip()} | {value.strip()}")
        else:
            lines.append(line)
    return "\n".join(lines)


def extract_numbers(metrics_text):
    """Best-effort cFTTH/cMobile extraction from the metrics text -- the
    old deck's phrasing isn't consistent enough to guarantee this, treat
    it as a hint to verify, not ground truth."""
    cFTTH = cMobile = ""
    m = re.search(r"FTTH[^0-9]*?([\d.,]+)\s*[kK]", metrics_text)
    if m:
        cFTTH = m.group(1).replace(".", "").replace(",", ".")
    m = re.search(r"([\d.,]+)\s*[kK]\s*clientes?\s*de\s*m[oó]vil", metrics_text, re.IGNORECASE)
    if m:
        cMobile = m.group(1).replace(".", "").replace(",", ".")
    return cFTTH, cMobile


def parse_incident_group(group_shape, group_label):
    items = []
    collect_textboxes([group_shape], items)

    # The ticket textbox is a *sibling* of the inner "card" group (not
    # nested inside the 3-column content), so it's already in `items`
    # alongside everything else -- but it's the only one containing what
    # looks like a bare ticket code. Extract lines that look like tickets
    # (letters+digits, no spaces) instead of relying on position for this
    # one, since its exact offset shifts slide to slide.
    ticket_candidates = []
    other_items = []
    for x, y, text in items:
        lines = [l.strip() for l in text.split("\n") if l.strip()]
        if lines and all(TICKET_RE.fullmatch(l) for l in lines):
            ticket_candidates.extend(lines)
        else:
            other_items.append((x, y, text))

    items = other_items
    exact = {t: (x, y) for x, y, t in items}

    title = None
    for x, y, text in items:
        if x < 1.5 and y < 2.0 and text not in LABEL_WORDS:
            title = text
            break

    fecha = nearest([(x, y, t) for x, y, t in items if 9.0 < x < 10.0 and y > 1.4 and t not in LABEL_WORDS], 9.28)
    duracion = nearest([(x, y, t) for x, y, t in items if x > 10.5 and y > 1.4 and t not in LABEL_WORDS], 11.12)

    body = [(x, y, t) for x, y, t in items if y > 3.8 and t not in LABEL_WORDS]
    impacto = nearest(body, 0.5, tolerance=1.0)
    causa = nearest(body, 4.0, tolerance=1.0)
    solucion = nearest(body, 8.1, tolerance=1.0)

    metrics_raw = impacto or ""
    cFTTH, cMobile = extract_numbers(metrics_raw)

    flags = []
    combined_lower = " ".join(filter(None, [title, causa, solucion, duracion])).lower()
    if "xxxx" in combined_lower or (duracion or "").strip() in ("", "h m"):
        flags.append("datos incompletos (placeholder xxxx/vacío) en el PPT original")

    return {
        "group": group_label,
        "severity": "SL2",
        "category": "",
        "system": "",
        "title": title or "",
        "ticket": " / ".join(ticket_candidates),
        "date": fecha or "",
        "duration": duracion or "",
        "impact": "",
        "metrics": metrics_to_pipe_format(metrics_raw),
        "cause": causa or "",
        "solution": solucion or "",
        "actionPoints": "",
        "cFTTH": cFTTH,
        "cMobile": cMobile,
        "brands": "",
        "ministry": False,
        "platform": False,
        "externalOrigin": False,
        "_flags": flags,
    }


def find_incident_groups(slide_shapes):
    """Top-level incident groups are direct children of the slide (not
    nested inside another group) whose text, once flattened, contains a
    ticket-shaped line -- this is more robust than assuming a fixed
    shape count, since some incidents have 1 ticket and some have 2."""
    groups = []
    for shape in slide_shapes:
        if shape.shape_type != 6:
            continue
        items = []
        collect_textboxes([shape], items)
        has_ticket = any(TICKET_RE.fullmatch(l.strip())
                          for _, _, t in items for l in t.split("\n") if l.strip())
        if has_ticket:
            groups.append(shape)
    return groups


def parse_pptx(path):
    prs = Presentation(path)
    incidents = []
    warnings = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        title_shape = next((s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()), None)
        title_text = title_shape.text_frame.text.strip() if title_shape else ""
        group_label, known = slide_group(title_text)
        if not known:
            warnings.append(f"Slide {slide_idx}: grupo '{group_label}' no es una convención existente en la app, revisar.")

        for grp in find_incident_groups(slide.shapes):
            inc = parse_incident_group(grp, group_label)
            flags = inc.pop("_flags")
            if not inc["title"]:
                warnings.append(f"Slide {slide_idx}: no se pudo extraer el título de una incidencia (ticket={inc['ticket']!r}).")
            for f in flags:
                warnings.append(f"Slide {slide_idx} / {inc['ticket'] or inc['title'][:40]}: {f}")
            incidents.append(inc)

    return incidents, warnings


def infer_year_week(path, year, week):
    if year and week:
        return year, week
    m = re.search(r"(\d{4})W(\d{2})", path)
    if m:
        return year or int(m.group(1)), week or int(m.group(2))
    return year, week


def infer_range(incidents):
    dates = []
    for inc in incidents:
        m = re.match(r"(\d{2})/(\d{2})/(\d{4})", inc["date"])
        if m:
            dates.append((m.group(3), m.group(2), m.group(1)))  # y, m, d for sorting
    if not dates:
        return ""
    dates.sort()
    d0, d1 = dates[0], dates[-1]
    return f"{d0[2]}-{d1[2]} de {['ene','feb','mar','abr','may','jun','jul','ago','sep','oct','nov','dic'][int(d1[1]) - 1]}"


def main():
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("pptx")
    ap.add_argument("--year", type=int, default=None)
    ap.add_argument("--week", type=int, default=None)
    ap.add_argument("--range", dest="range_", default=None)
    ap.add_argument("--dept", default="Customer & Service Operations")
    ap.add_argument("-o", "--output", default=None)
    ap.add_argument("--post-url", default=None, help="e.g. http://localhost:8000 -- POSTs the report directly instead of/as well as writing JSON")
    args = ap.parse_args()

    incidents, warnings = parse_pptx(args.pptx)
    year, week = infer_year_week(args.pptx, args.year, args.week)
    if not year or not week:
        print("ERROR: no se pudo inferir year/week del nombre de fichero, pásalos con --year/--week", file=sys.stderr)
        sys.exit(1)

    report = {
        "year": year,
        "week": week,
        "range": args.range_ or infer_range(incidents),
        "dept": args.dept,
        "incidents": incidents,
        "status": "draft",
    }

    print(f"Extraídas {len(incidents)} incidencias de {args.pptx}")
    print(f"year={year} week={week} range={report['range']!r} dept={report['dept']!r}\n")

    if warnings:
        print(f"⚠️  {len(warnings)} avisos -- revisar a mano tras importar:")
        for w in warnings:
            print(f"  - {w}")
        print()
    print("Campos que ESTE PPT no registraba y quedan con su valor por defecto en TODAS las incidencias:")
    print("  - severity (default SL2), brands (vacío), ministry/platform/externalOrigin (False), actionPoints (vacío), category/system (vacío, todo va en 'title')")
    print()

    if args.output:
        with open(args.output, "w", encoding="utf-8") as f:
            json.dump(report, f, ensure_ascii=False, indent=2)
        print(f"JSON escrito en {args.output}")

    if args.post_url:
        data = json.dumps(report).encode("utf-8")
        req = urllib.request.Request(
            f"{args.post_url}/api/reports", data=data,
            headers={"Content-Type": "application/json"}, method="POST",
        )
        try:
            with urllib.request.urlopen(req) as resp:
                print(f"POST {args.post_url}/api/reports -> {resp.status}")
                print(resp.read().decode("utf-8")[:500])
        except urllib.error.HTTPError as e:
            print(f"POST falló: {e.code} {e.read().decode('utf-8')}", file=sys.stderr)
            sys.exit(1)


if __name__ == "__main__":
    main()
